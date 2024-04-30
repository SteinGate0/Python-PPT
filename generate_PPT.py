import json5
from pptx import Presentation
from datetime import datetime
import re
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
import random

def parse_content(content):
    try:
        match = re.search(r"(\{.*\})", content, re.DOTALL)
        if match:
            content = match.groups()[0]
            print("match content:")
            print(content)
        return json5.loads(content.strip())
    except Exception as e:
        print(f"The response is not a valid JSON format: {e}")
        print("I'm a PPT assistant, your PPT generate failed, please retry later..")
        raise Exception("The LLM return invalid result, please retry later..")
        exit(1)


ppt_content = '''{
    'title': '养生',
    'pages': [
        {
            'title': '标题1',
            'content': [
                {'title': '标题1.1', 'description': '详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
                {'title': '标题1.2', 'description': '详细描述标题1.2的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
                {'title': '标题1.3', 'description': '详细描述标题1.3的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
                {'title': '标题1.3', 'description': '详细描述标题1.3的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
            ]
        },
        {
            'title': '标题2',
            'content': [
                {'title': '标题2.1', 'description': '详细描述标题2.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
                {'title': '标题2.2', 'description': '详细描述标题2.2的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
                {'title': '标题2.3', 'description': '详细描述标题2.2的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
            ]
        },
        {
            'title': '标题3',
            'content': [
                {'title': '标题3.1', 'description': '详细描述标题3.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
                {'title': '标题3.2', 'description': '详细描述标题3.2的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
                {'title': '标题3.3', 'description': '详细描述标题3.3的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
            ]
        },
        {
            'title': '标题4',
            'content': [
                {'title': '标题4.1', 'description': '详细描述标题4.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
                {'title': '标题4.2', 'description': '详细描述标题4.2的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
                {'title': '标题4.3', 'description': '详细描述标题4.3的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容详细描述标题1.1的内容'},
            ]
        }
    ]
}'''


def font_style(font, size, bold=False, R=0, G=0, B=255):
    font.name = '微软雅黑'  # 设置字体名称，例如宋体、微软雅黑等
    font.size = Pt(size)  # 设置字体大小，单位为磅
    font.bold = bold  # 是否加粗
    font.color.rgb = RGBColor(R, G, B)  # 设置字体颜色，RGBColor(r, g, b)


def content_PPT1(page, ppt):
    content = ppt.slides.add_slide(ppt.slide_layouts[1])
    text_box1 = content.shapes.add_textbox(Cm(2.69), Cm(2.19), Cm(26.03), Cm(1.98))
    tf1 = text_box1.text_frame
    tf1.text = page['title']
    # tf1.margin_bottom = tf1.margin_top = Cm(0.13)  # 下边距
    # tf1.margin_left = tf1.margin_right= Cm(0.25)   # 左左距
    tf1.vertical_anchor = MSO_ANCHOR.BOTTOM
    font1 = tf1.paragraphs[0].font
    font_style(font1, 30, True, 0, 0, 0)

    for i, item in enumerate(page['content']):
        left = i * 8.8 + 3.91  # 0 0   1 8.8   2 2*8.8
        text_box2 = content.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(5.91), Cm(8.45), Cm(10.08))
        tf2 = text_box2.text_frame
        tf2.text = item['title'] + "\n"
        font2 = tf2.paragraphs[0].font
        font_style(font2, 21, True, 0, 0, 128)
        # -----------
        p2 = tf2.add_paragraph()
        run2 = p2.add_run()
        run2.text = item['description']
        font3 = run2.font
        font_style(font3, 15.8, False, 0, 0, 0)

        # 设置文本框格式
        text_box2.fill.solid()  # 填充为纯色
        text_box2.fill.fore_color.rgb = RGBColor(240, 248, 255)  # 设置填充颜色，RGBColor(r, g, b)
        text_box2.line.color.rgb = RGBColor(0, 0, 0)  # 设置边框颜色，RGBColor(r, g, b)
        text_box2.line.width = Pt(1)  # 设置边框宽度，单位为磅


def content_PPT2(page, ppt):
    content = ppt.slides.add_slide(ppt.slide_layouts[1])
    text_box1 = content.shapes.add_textbox(Cm(2.69), Cm(2.19), Cm(26.03), Cm(1.98))
    tf1 = text_box1.text_frame
    tf1.text = page['title']
    tf1.vertical_anchor = MSO_ANCHOR.BOTTOM
    font1 = tf1.paragraphs[0].font
    font_style(font1, 30, True, 0, 0, 0)
    # 135,206,250
    for i, item in enumerate(page['content']):
        left = i * 7.98 + 4.93  # 0 0   7.98    2 7.98
        text_box2 = content.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(6.01), Cm(7.04), Cm(1.8))
        text_box2.fill.solid()  # 填充为纯色
        text_box2.fill.fore_color.rgb = RGBColor(240, 248, 255)  # 设置填充颜色，RGBColor(r, g, b)
        text_box2.line.color.rgb = RGBColor(30, 144, 255)  # 设置边框颜色，RGBColor(r, g, b)

        tf2 = text_box2.text_frame
        tf2.text = item['title']
        font2 = tf2.paragraphs[0].font
        font_style(font2, 21, True, 0, 0, 128)
        # -----------
        left1 = i * 7.98 + 4.93
        text_box = content.shapes.add_textbox(Cm(left1), Cm(8.16), Cm(7.62), Cm(4.26))
        tf3 = text_box.text_frame
        tf3.text = item['description']
        tf3.word_wrap = True  # 框中的文字自动换行
        font3 = tf3.paragraphs[0].font
        font_style(font3, 15.8, False, 0, 0, 0)


def content_PPT3(page, ppt):
    content = ppt.slides.add_slide(ppt.slide_layouts[1])
    text_box1 = content.shapes.add_textbox(Cm(2.69), Cm(2.19), Cm(26.03), Cm(1.98))
    tf1 = text_box1.text_frame
    tf1.text = page['title']
    tf1.vertical_anchor = MSO_ANCHOR.BOTTOM
    font1 = tf1.paragraphs[0].font
    font_style(font1, 30, True, 0, 0, 0)

    for i, item in enumerate(page['content']):
        top = i * 4.26 + 5.22  # 0 0   1 8.8   2 2*8.8
        text_box2 = content.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(3.53), Cm(top), Cm(1), Cm(1.03))
        text_box2.fill.solid()  # 填充为纯色
        text_box2.fill.fore_color.rgb = RGBColor(0, 0, 255)  # 设置填充颜色，RGBColor(r, g, b)

        tf2 = text_box2.text_frame
        tf2.text = str(i+1)
        font2 = tf2.paragraphs[0].font
        font_style(font2, 18, False, 255, 255, 255)
        # -----------
        text_box3 = content.shapes.add_textbox(Cm(4.93), Cm(top), Cm(22.27), Cm(3.91))
        tf3 = text_box3.text_frame
        tf3.word_wrap = True  # 框中的文字自动换行
        tf3.text = item['title'] + "\n"
        font3 = tf3.paragraphs[0].font
        font_style(font3, 21, True, 0, 0, 128)
        #  ----------------

        p3 = tf3.add_paragraph()
        run3 = p3.add_run()
        run3.text = item['description']
        font4 = run3.font
        font_style(font4, 15.8, False, 0, 0, 0)



def content_PPT4(page, ppt):
    content = ppt.slides.add_slide(ppt.slide_layouts[1])
    text_box1 = content.shapes.add_textbox(Cm(2.69), Cm(2.19), Cm(26.03), Cm(1.98))
    tf1 = text_box1.text_frame
    tf1.text = page['title']
    tf1.vertical_anchor = MSO_ANCHOR.BOTTOM
    font1 = tf1.paragraphs[0].font
    font_style(font1, 30, True, 0, 0, 0)

    for i, item in enumerate(page['content']):
        left = i * 7.98 + 4.93  # 0 0   1 8.8   2 2*8.8
        # -----------
        text_box2 = content.shapes.add_textbox(Cm(left), Cm(6.05), Cm(7.62), Cm(5.91))
        tf2 = text_box2.text_frame
        tf2.word_wrap = True  # 框中的文字自动换行
        tf2.text = item['title'] + "\n"
        font3 = tf2.paragraphs[0].font
        font_style(font3, 21, True, 0, 0, 128)
        #  ----------------
        p2 = tf2.add_paragraph()
        run2 = p2.add_run()
        run2.text = item['description']
        font2 = run2.font
        font_style(font2, 15.8, False, 0, 0, 0)

def content_PPT5(page, ppt):
    content = ppt.slides.add_slide(ppt.slide_layouts[1])
    text_box1 = content.shapes.add_textbox(Cm(2.69), Cm(2.19), Cm(26.03), Cm(1.98))
    tf1 = text_box1.text_frame
    tf1.text = page['title']
    tf1.vertical_anchor = MSO_ANCHOR.BOTTOM
    font1 = tf1.paragraphs[0].font
    font_style(font1, 30, True, 0, 0, 0)

    for i, item in enumerate(page['content']):
        top = i * 2.32 + 4.75  # 0 0   1 8.8   2 2*8.8
        text_box2 = content.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(16.1), Cm(top), Cm(1), Cm(1.03))
        text_box2.fill.solid()  # 填充为纯色
        text_box2.fill.fore_color.rgb = RGBColor(0, 0, 255)  # 设置填充颜色，RGBColor(r, g, b)

        tf2 = text_box2.text_frame
        tf2.text = str(i+1)
        font2 = tf2.paragraphs[0].font
        font_style(font2, 18, False, 255, 255, 255)
        # -----------
        if i % 2 ==0 :
            text_box3 = content.shapes.add_textbox(Cm(17.91), Cm(top), Cm(12.63), Cm(4.15))
            tf3 = text_box3.text_frame
            tf3.word_wrap = True  # 框中的文字自动换行
            tf3.text = item['title']
            font3 = tf3.paragraphs[0].font
            font_style(font3, 21, True, 0, 0, 128)
            #  ----------------
            p3 = tf3.add_paragraph()
            run3 = p3.add_run()
            run3.text = item['description']
            font4 = run3.font
            font_style(font4, 15.8, False, 0, 0, 0)
            #-----------样式
            top1 = i/2 * 4.63 + 5.25  # 0 0   2  5.03
            shape = content.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(17.13), Cm(top1), Cm(0.86), Cm(0.11))
            shape.fill.solid()  # 填充为纯色
            shape.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 设置填充颜色，RGBColor(r, g, b)

        else:
            text_box3 = content.shapes.add_textbox(Cm(2.69), Cm(top), Cm(12.63), Cm(4.15))
            tf3 = text_box3.text_frame
            tf3.word_wrap = True  # 框中的文字自动换行

            tf3.text = item['title']
            tf3.paragraphs[0].alignment = PP_ALIGN.RIGHT
            font3 = tf3.paragraphs[0].font
            font_style(font3, 21, True, 0, 0, 128)
            #  ----------------
            p3 = tf3.add_paragraph()
            run3 = p3.add_run()
            p3.alignment = PP_ALIGN.RIGHT
            run3.text = item['description']
            font4 = run3.font
            font_style(font4, 15.8, False, 0, 0, 0)
            #  样式
            top2 = (i-1)/2 * 4.63 + 7.57  # 1 0   3  5.03
            shape = content.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(15.19), Cm(top2), Cm(0.86), Cm(0.11))
            shape.fill.solid()  # 填充为纯色
            shape.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 设置填充颜色，RGBColor(r, g, b)
        top3 = i  * 2.32 + 5.78
        shape = content.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(16.54), Cm(top3), Cm(0.11), Cm(1.29))
        shape.fill.solid()  # 填充为纯色
        shape.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 设置填充颜色，RGBColor(r, g, b)

def generate_ppt_file(topic, ppt_content):
    ppt = Presentation("Hicon-st.pptx")
    ppt_content = parse_content(ppt_content)
    print(ppt_content['title'])
    # PPT首页
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # title&subtitle layout
    print("首页")
    for shape in slide.placeholders:  # 获取这一页所有的占位符
        phf = shape.placeholder_format
        print(f'{phf.idx}--{shape.name}--{phf.type}')  # id号--占位符形状名称-占位符的类型
    # 首页
    slide.placeholders[0].text = ppt_content['title']
    slide.placeholders[1].text = "Hicon-ChatOps"
    slide.placeholders[10].text = "日期：" + str(datetime.today().date())
    slide.placeholders[11].text = "主讲人："

    catalogue = ppt.slides.add_slide(ppt.slide_layouts[-2])
    # 目录
    for i, page in enumerate(ppt_content['pages']):
        if i % 2 == 0:
            left = 6.5
            top = 6.93 + i / 2 * 1.4  # 2   1.4    4 2.8     6
            text_box = catalogue.shapes.add_textbox(Cm(left), Cm(top), Cm(11.54), Cm(1.4))  # left top
            tf = text_box.text_frame
            tf.text = page['title']
            font = tf.paragraphs[0].font
            font.name = '微软雅黑'  # 设置字体名称，例如宋体、微软雅黑等
            font.size = Pt(24)  # 设置字体大小，单位为磅
            font.bold = True  # 是否加粗
            font.color.rgb = RGBColor(0, 0, 255)  # 设置字体颜色，RGBColor(r, g, b)
        else:
            left = 6.5 + 11.54
            top = 6.93 + (i - 1) / 2 * 1.4  # 1 0    3  1.4    5  2.8
            text_box = catalogue.shapes.add_textbox(Cm(left), Cm(top), Cm(11.54), Cm(1.4))  # left top
            tf = text_box.text_frame
            tf.text = page['title']
            font = tf.paragraphs[0].font
            font.name = '微软雅黑'  # 设置字体名称，例如宋体、微软雅黑等
            font.size = Pt(24)  # 设置字体大小，单位为磅
            font.bold = True  # 是否加粗
            font.color.rgb = RGBColor(0, 0, 255)  # 设置字体颜色，RGBColor(r, g, b)

    for i, page in enumerate(ppt_content['pages']):
        guide_slide = ppt.slides.add_slide(ppt.slide_layouts[-3])
        guide_slide.placeholders[1].text = page['title']
        guide_slide.placeholders[10].text = "0" + str(i + 1)
        content_PPT = [content_PPT1, content_PPT2, content_PPT3, content_PPT4, content_PPT5]
        content_PPT = random.choice(content_PPT)
        content_PPT(page, ppt)

    slide = ppt.slides.add_slide(ppt.slide_layouts[-1])
    ppt.save('%s.pptx' % topic)


generate_ppt_file("example", ppt_content)
