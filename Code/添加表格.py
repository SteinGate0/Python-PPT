from pptx import Presentation
from pptx.util import Cm

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6] #使用第7个空白版式创建新的幻灯片
slide = prs.slides.add_slide(blank_slide_layout)

rows,cols = 4,2
left = top =Cm(5)
width = Cm(18)
height = Cm(3)

table = slide.shapes.add_table(rows,cols,left,top,width,height).table
# 可以修改列宽、行高
table.columns[0].width = Cm(6)
table.columns[1].width = Cm(4)
table.rows[0].height =Cm(2)

data = [
    ['姓名','成绩'],
    ['李雷',99],
    ['韩梅梅', 92],
    ['马东梅', 92],
]
for row in range(rows):
    for col in range(cols):
        table.cell(row,col).text =str(data[row][col])

prs.save('添加表格.pptx')