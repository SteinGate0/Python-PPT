from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

SLD_LAYOUT_TITLE_AND_CONTENT = 6
prs = Presentation()
slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
slide = prs.slides.add_slide(slide_layout)

# 支持的形状可以看：
# https://docs.microsoft.com/zh-cn/office/vba/api/Office.MsoAutoShapeType
shapes = slide.shapes
left, top, width, height = Inches(1), Inches(3), Inches(1.8), Inches(1)
shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
shape.text ='第一步'

for n in range(2,6):
    left=left+width-Inches(-0.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON,left,top,width,height)
    shape.text = f'第{n}步'
prs.save('添加形状.pptx')
