from pptx import Presentation

ppt = Presentation()
print("内容页")
content_slide = ppt.slides.add_slide(ppt.slide_layouts[2])
for shape in content_slide.placeholders:  # 获取这一页所有的占位符
  phf = shape.placeholder_format
  print(f'{phf.idx}--{shape.name}--{phf.type}')  # id号--占位符形状名称-占位符的类型