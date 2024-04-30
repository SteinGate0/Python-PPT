from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)

left = top = width = height = Cm(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
tf.text = "这是一段文本框里面的文字"
# -------------------------------------- #
tf.margin_bottom = Cm(0.1) # 下边距
tf.margin_left = 0 # 坐左距
tf.vertical_anchor = MSO_ANCHOR.BOTTOM
tf.word_wrap = True # 框中的文字自动换行
# -------------------------------------- #
fill = text_box.fill
fill.solid()
# 使用之前一定要导入RGBColor这个库
fill.fore_color.rgb = RGBColor(247, 150, 70)
# -------------------------------------- #
line = text_box.line
line.color.rgb = RGBColor(255, 0, 0)
line.width = Cm(0.3)

prs.save("文本框边框样式调整.pptx")
