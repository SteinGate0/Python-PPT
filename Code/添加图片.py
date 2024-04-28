from pptx import Presentation
from pptx.util import Cm

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Cm(3)
pic = slide.shapes.add_picture('img/animal.png',left,top)

left = height = Cm(10)
pic = slide.shapes.add_picture('img/animal.png',left,top, height=height)

prs.save('添加图片.pptx')