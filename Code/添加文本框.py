from pptx import Presentation
from pptx.util import Cm,Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height =Cm(3)  # # left，top为相对位置，width，height为文本框大小。满足条件顺序是左>上>右>下
text_box = slide.shapes.add_textbox(left,top,width,height)
tf = text_box.text_frame
tf.text = "这是一段文本框里的文字"

p = tf.add_paragraph()
p.text = "这是第二段文字，加粗，字号40"
p.font.bold = True
p.font.size = Pt(40)

prs.save('添加文本框.pptx')
