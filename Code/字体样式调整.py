from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)

left = top = width = height = Cm(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
# ---------------------------- #
p = tf.add_paragraph()
p.text = "这是第二段文字"
p.alignment = PP_ALIGN.LEFT
# ------------------------------------- #
p.font.bold = True
p.font.name = "宋体"
p.font.color.rgb = RGBColor(247, 150, 70)
p.font.size = Pt(30)

prs.save("字体样式调整.pptx")
