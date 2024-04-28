from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

prs = Presentation()
# 使用第七个版式
black_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(black_slide_layout)

left = top = width = height = Cm(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
# ---------------------------- #
p = tf.add_paragraph()
p.text = "这是第二段文字"
p.alignment = PP_ALIGN.LEFT  # 左对齐
#p.add_run().text = "新文字"

prs.save("段落对其调整.pptx")
